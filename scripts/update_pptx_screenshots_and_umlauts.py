"""
1. Replace old screenshots on cluster slides with new dashboard screenshots
2. Fix all Umlauts: ue->ü, ae->ä, oe->ö, Ue->Ü, Ae->Ä, Oe->Ö
"""
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
import re

pptx_path = r"C:\Users\bensc\OneDrive\HWR\6. Semester\Bachelor\07_Prototypen\mvp_v3.0_new\TI-Radar_MVP_v3.0_Praesentation.pptx"
screenshot_dir = r"C:\Users\bensc\OneDrive\HWR\6. Semester\Bachelor\07_Prototypen\mvp_v3.0_new\screenshots"

prs = Presentation(pptx_path)
sw = prs.slide_width
sh = prs.slide_height

# ============================================================
# PART 1: Find cluster slides and replace/add screenshots
# ============================================================

# Map slide titles to new screenshots
SCREENSHOT_MAP = {
    "CLUSTER 1": f"{screenshot_dir}\\02_skurve_reife.png",
    "CLUSTER 2": f"{screenshot_dir}\\03_marktakteure.png",
    "CLUSTER 3": f"{screenshot_dir}\\04_foerderung.png",
    "CLUSTER 4": f"{screenshot_dir}\\05_geographie.png",
}

# Also update Dashboard slide
DASHBOARD_SCREENSHOT = f"{screenshot_dir}\\01_dashboard_overview.png"

screenshots_replaced = 0

for slide in prs.slides:
    # Get slide title
    title = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                title = t
                break

    # Check if this slide needs a screenshot
    target_img = None
    if "DASHBOARD" in title:
        target_img = DASHBOARD_SCREENSHOT
    else:
        for key, img_path in SCREENSHOT_MAP.items():
            if key in title:
                target_img = img_path
                break

    if target_img is None:
        continue

    # Find existing non-background picture shapes and remove them
    pics_to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == 13:  # Picture
            # Skip background (full slide size at 0,0)
            if shape.left == 0 and shape.top == 0 and shape.width == sw:
                continue
            # Skip small images (icons, ERM thumbnails)
            if shape.width < Emu(3000000):  # < ~8cm
                continue
            pics_to_remove.append(shape)

    for pic in pics_to_remove:
        sp_tree = slide.shapes._spTree
        sp_tree.remove(pic._element)

    # Add new screenshot
    # Position: centered, below title area, fitting most of the slide
    img_left = Emu(500000)
    img_top = Emu(1000000)
    img_width = sw - Emu(1000000)
    img_height = sh - Emu(1300000)

    slide.shapes.add_picture(target_img, img_left, img_top, img_width, img_height)
    screenshots_replaced += 1
    print(f"  Screenshot: {title[:40]} <- {target_img.split(chr(92))[-1]}")

print(f"\n{screenshots_replaced} screenshots replaced/added.")

# ============================================================
# PART 2: Fix Umlauts in all text
# ============================================================

# Words that should NOT be converted (false positives)
EXCEPTIONS = {
    "Bluetooth", "Queue", "Blue", "True", "Value", "Issue",
    "Continue", "Rescue", "Revenue", "Venue", "Tissue",
    "Tableau", "Bureau", "Trousseau", "Plateau",
    "Route", "Compute", "Execute", "Module", "Node",
    "Unique", "Technique", "Critique", "Fatigue",
    "Ogue", "Vogue", "Rogue", "Analogue",
    "Manuel", "Joel", "Samuel", "Noel", "Israel",
    "Michael", "Raphael", "Fuel", "Cruel", "Duel",
    "Does", "Goes", "Shoes", "Toes", "Foes", "Hoes",
    "Poet", "Poem", "Poetry", "Aloe", "Canoe", "Oboe",
    "Phoenix", "Coefficient", "Hoe", "Joe", "Doe", "Roe",
    "Coe", "Zoe",
}
EXCEPTIONS_LOWER = {w.lower() for w in EXCEPTIONS}


def fix_umlauts(text):
    """Replace ae/oe/ue with ä/ö/ü in German words, avoiding English false positives."""
    if not text:
        return text

    result = text

    # Direct replacements for common German patterns
    replacements = [
        # ue -> ü
        ("ueberblick", "Überblick"),
        ("Ueberblick", "Überblick"),
        ("UEBERBLICK", "ÜBERBLICK"),
        ("ueberblick", "überblick"),
        ("uebersicht", "Übersicht"),
        ("ueber", "über"),
        ("Ueber", "Über"),
        ("fuer", "für"),
        ("Fuer", "Für"),
        ("Einfuehrung", "Einführung"),
        ("Durchfuehrung", "Durchführung"),
        ("Ausfuehrung", "Ausführung"),
        ("Verfuegbar", "Verfügbar"),
        ("verfuegbar", "verfügbar"),
        ("Pruefung", "Prüfung"),
        ("gruende", "gründe"),
        ("Gruende", "Gründe"),
        ("Stueck", "Stück"),
        ("Unterstuetzung", "Unterstützung"),
        ("unterstuetzt", "unterstützt"),
        ("zurueck", "zurück"),
        ("Zurueck", "Zurück"),
        ("natuerlich", "natürlich"),
        ("kuenstlich", "künstlich"),
        ("Schluessel", "Schlüssel"),
        ("SCHLUESSEL", "SCHLÜSSEL"),
        ("Schluesselmerkmal", "Schlüsselmerkmal"),
        ("SCHLUESSELMERKMALE", "SCHLÜSSELMERKMALE"),
        ("Ausfuellen", "Ausfüllen"),
        ("Spuerbar", "Spürbar"),
        ("Ergebnisse und naechste", "Ergebnisse und nächste"),
        ("naechste", "nächste"),
        ("Naechste", "Nächste"),

        # ae -> ä
        ("Aktivitaet", "Aktivität"),
        ("Qualitaet", "Qualität"),
        ("QUALITAET", "QUALITÄT"),
        ("Stabilitaet", "Stabilität"),
        ("Universitaet", "Universität"),
        ("Intensitaet", "Intensität"),
        ("Komplexitaet", "Komplexität"),
        ("Kreativitaet", "Kreativität"),
        ("Prioritaet", "Priorität"),
        ("Realitaet", "Realität"),
        ("Skaerfe", "Schärfe"),
        ("Staerke", "Stärke"),
        ("staerke", "stärke"),
        ("Schaerfe", "Schärfe"),
        ("spaeter", "später"),
        ("Spaeter", "Später"),
        ("waehrend", "während"),
        ("Waehrend", "Während"),
        ("ungefaehr", "ungefähr"),
        ("Laender", "Länder"),
        ("laender", "länder"),
        ("LAENDER", "LÄNDER"),
        ("Aenderung", "Änderung"),
        ("aendern", "ändern"),
        ("Geraet", "Gerät"),
        ("Geraete", "Geräte"),
        ("Datensaetze", "Datensätze"),
        ("Datensaetzen", "Datensätzen"),
        ("Einsaetze", "Einsätze"),
        ("spaeter", "später"),
        ("Praesentation", "Präsentation"),
        ("praesent", "präsent"),
        ("Uebersetzung", "Übersetzung"),
        ("Erlaeuterung", "Erläuterung"),
        ("erlaeutern", "erläutern"),

        # oe -> ö
        ("Foerderung", "Förderung"),
        ("foerderung", "förderung"),
        ("FOERDERUNG", "FÖRDERUNG"),
        ("Loesung", "Lösung"),
        ("loesung", "lösung"),
        ("LOESUNG", "LÖSUNG"),
        ("moeglich", "möglich"),
        ("Moeglich", "Möglich"),
        ("erhoehen", "erhöhen"),
        ("Erhoehung", "Erhöhung"),
        ("gehoert", "gehört"),
        ("Behoerde", "Behörde"),
        ("Groesse", "Größe"),
        ("groesste", "größte"),
        ("Groesste", "Größte"),
        ("Stoerung", "Störung"),

        # Common compound words
        ("zeitaufwaendig", "zeitaufwändig"),
        ("fehleranfaellig", "fehleranfällig"),
        ("Technologiefeldern", "Technologiefeldern"),  # no change needed
        ("Quellenangaben", "Quellenangaben"),  # no change needed

        # Specific presentation text fixes
        ("Forschungsfoerderung", "Forschungsförderung"),
        ("EU-Forschungsfoerderung", "EU-Forschungsförderung"),
        ("Technologiebewertung", "Technologiebewertung"),  # correct already
        ("Technologiekonvergenz", "Technologiekonvergenz"),  # correct already
        ("Technologieprognosen", "Technologieprognosen"),  # correct already
        ("Technologieanalyse", "Technologieanalyse"),  # correct already
        ("gestuetzt", "gestützt"),
        ("datengestuetzt", "datengestützt"),
        ("KI-gestuetzt", "KI-gestützt"),
        ("Produktionsbetrieb", "Produktionsbetrieb"),  # correct
        ("Verschluesselung", "Verschlüsselung"),
        ("verschluesselt", "verschlüsselt"),
        ("verschluesselte", "verschlüsselte"),
    ]

    for old, new in replacements:
        if old in result:
            result = result.replace(old, new)

    return result


umlaut_fixes = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            txBody = shape.text_frame._txBody
            for t_elem in txBody.findall(".//" + qn("a:t")):
                if t_elem.text:
                    fixed = fix_umlauts(t_elem.text)
                    if fixed != t_elem.text:
                        old = t_elem.text
                        t_elem.text = fixed
                        umlaut_fixes += 1
                        if len(old) < 80:
                            try:
                                print(f"  Umlaut: '{old}' -> '{fixed}'")
                            except UnicodeEncodeError:
                                print(f"  Umlaut: (fixed, console encoding issue)")

print(f"\n{umlaut_fixes} Umlaut-Korrekturen.")

prs.save(pptx_path)
print(f"Saved: {pptx_path}")
