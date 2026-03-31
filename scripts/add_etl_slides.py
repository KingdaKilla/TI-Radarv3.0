"""Add 3 ETL Pipeline slides after the Datenkonzept block (after Folie 12)."""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN
from lxml import etree

pptx_path = r"C:\Users\bensc\OneDrive\HWR\6. Semester\Bachelor\07_Prototypen\mvp_v3.0_new\TI-Radar_MVP_v3.0_Praesentation.pptx"
bg_image = r"C:\Users\bensc\OneDrive\HWR\6. Semester\Bachelor\07_Prototypen\mvp_v3.0_new\figures\slide_background_dark_v2.png"

prs = Presentation(pptx_path)
sw = prs.slide_width
sh = prs.slide_height
blank = prs.slide_layouts[6]
M = 457200  # margin


def setup_slide(prs):
    slide = prs.slides.add_slide(blank)
    pic = slide.shapes.add_picture(bg_image, 0, 0, sw, sh)
    sp = slide.shapes._spTree
    el = pic._element
    sp.remove(el)
    sp.insert(2, el)
    bg = slide.background
    bgPr = bg._element.find(qn("p:bgPr"))
    if bgPr is not None:
        sf = bgPr.find(qn("a:solidFill"))
        if sf is not None:
            bgPr.remove(sf)
            etree.SubElement(bgPr, qn("a:noFill"))
    return slide


def txt(slide, left, top, width, height, text, size, color="B8CAD4",
        bold=False, align="left"):
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    pPr = p._pPr
    if pPr is None:
        pPr = etree.SubElement(p._p, qn("a:pPr"))
    dr = etree.SubElement(pPr, qn("a:defRPr"))
    dr.set("sz", str(int(size * 100)))
    dr.set("b", "1" if bold else "0")
    sf = etree.SubElement(dr, qn("a:solidFill"))
    sr = etree.SubElement(sf, qn("a:srgbClr"))
    sr.set("val", color)
    la = etree.SubElement(dr, qn("a:latin"))
    la.set("typeface", "Segoe UI")
    p.add_run().text = text
    return tb


def box(slide, left, top, width, height, fill="162138", alpha="80000",
        border="1B4D6E"):
    shape = slide.shapes.add_shape(1, Emu(left), Emu(top), Emu(width), Emu(height))
    spPr = shape._element.find(qn("p:spPr"))
    for c in list(spPr):
        if c.tag in (qn("a:solidFill"), qn("a:noFill")):
            spPr.remove(c)
    sf = etree.SubElement(spPr, qn("a:solidFill"))
    sr = etree.SubElement(sf, qn("a:srgbClr"))
    sr.set("val", fill)
    a = etree.SubElement(sr, qn("a:alpha"))
    a.set("val", alpha)
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    ln.set("w", str(int(Pt(0.75))))
    for c in list(ln):
        if c.tag in (qn("a:solidFill"), qn("a:noFill")):
            ln.remove(c)
    lf = etree.SubElement(ln, qn("a:solidFill"))
    ls = etree.SubElement(lf, qn("a:srgbClr"))
    ls.set("val", border)
    la = etree.SubElement(ls, qn("a:alpha"))
    la.set("val", "60000")
    pg = spPr.find(qn("a:prstGeom"))
    if pg is not None:
        pg.set("prst", "roundRect")
    return shape


def accent_bar(slide, x, y, height, color):
    bar = slide.shapes.add_shape(1, Emu(x), Emu(y), Emu(60000), Emu(height))
    spPr = bar._element.find(qn("p:spPr"))
    for c in list(spPr):
        if c.tag in (qn("a:solidFill"), qn("a:noFill")):
            spPr.remove(c)
    sf = etree.SubElement(spPr, qn("a:solidFill"))
    sr = etree.SubElement(sf, qn("a:srgbClr"))
    sr.set("val", color)
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    for c in list(ln):
        ln.remove(c)
    etree.SubElement(ln, qn("a:noFill"))


def arrow_box(slide, x, y, w, h, text, color):
    """Small colored box with text, used for flow arrows."""
    box(slide, x, y, w, h, fill=color, alpha="90000", border=color)
    txt(slide, x + 50000, y + 30000, w - 100000, h - 60000, text, 11,
        color="FFFFFF", bold=True, align="center")


# ================================================================
# FOLIE 1: ETL-PIPELINE — ÜBERBLICK
# ================================================================
s = setup_slide(prs)
txt(s, M, 150000, sw - 2 * M, 450000, "ETL-PIPELINE", 32,
    color="FFFFFF", bold=True, align="center")
txt(s, M, 480000, sw - 2 * M, 300000,
    "Datenquellen \u2192 Import \u2192 Transformation \u2192 Analyse", 16,
    color="8FA8B5", align="center")

# Flow: 4 stages left to right
stage_w = 2500000
stage_h = 4800000
gap = 250000
start_x = M + 200000
y = 1000000

stages = [
    ("DATENQUELLEN", "00BCC4", [
        ("EPO DOCDB", "154.8M Patente\n195 GB Nested XML-ZIPs"),
        ("CORDIS", "80.5K EU-Projekte\nJSON/CSV-Archive"),
        ("EuroSciVoc", "Wissenschafts-Taxonomie\nHierarchisch (220K)"),
        ("OpenAIRE / S2", "Publikationen + Zitationen\nAPI-Cache (7d / 30d TTL)"),
    ]),
    ("IMPORT", "FFB800", [
        ("Scheduler", "APScheduler Cron\nSonntag 02:00 UTC"),
        ("Bulk-Import", "asyncpg COPY-Protokoll\n10K Records/Batch"),
        ("Inkrementell", "import_log Tracking\nSkip bereits importiert"),
        ("Validierung", "CHECK Constraints\nTyp-Konvertierung"),
    ]),
    ("TRANSFORMATION", "E91E63", [
        ("Normalisierung", "CPC-Codes bereinigt\nApplicant-Namen unified"),
        ("Partitionierung", "Range-Partition by Year\nCo-Partitioning (3 Tab.)"),
        ("Entity Resolution", "pg_trgm Fuzzy-Matching\nGLEIF LEI-Lookup"),
        ("MV-Refresh", "9 Materialized Views\nCONCURRENTLY (non-blocking)"),
    ]),
    ("ANALYSE", "4CAF50", [
        ("13 UC-Services", "gRPC Read-Only\nPer-Service DB-Rollen"),
        ("Orchestrator", "Fan-Out asyncio.gather\nREST/JSON Gateway"),
        ("Dashboard", "Next.js 14 Frontend\nRecharts + Nivo + D3"),
        ("Export", "CSV + Analysis Cache\n24h TTL (JSONB)"),
    ]),
]

for idx, (title, color, items) in enumerate(stages):
    x = start_x + idx * (stage_w + gap)
    box(s, x, y, stage_w, stage_h)
    accent_bar(s, x, y, stage_h, color)

    txt(s, x + 150000, y + 100000, stage_w - 200000, 350000,
        title, 16, color=color, bold=True, align="center")

    for i, (item_title, item_desc) in enumerate(items):
        iy = y + 550000 + i * 1050000
        txt(s, x + 180000, iy, stage_w - 300000, 300000,
            item_title, 13, color="FFFFFF", bold=True)
        txt(s, x + 180000, iy + 300000, stage_w - 300000, 500000,
            item_desc, 11, color="8FA8B5")

    # Arrow between stages
    if idx < 3:
        ax = x + stage_w + 30000
        ay = y + stage_h // 2 - 150000
        txt(s, ax, ay, gap - 60000, 300000, "\u25B6", 20,
            color="78909C", align="center")

print("Added: ETL-Pipeline Überblick")

# ================================================================
# FOLIE 2: IMPORT-PIPELINE — DETAILS
# ================================================================
s = setup_slide(prs)
txt(s, M, 150000, sw - 2 * M, 450000, "IMPORT-PIPELINE", 32,
    color="FFFFFF", bold=True, align="center")

# Three import pipelines side by side
col_w = 3500000
col_gap = 300000
col_start_x = M + 200000
col_y = 800000
col_h = 5500000

pipelines = [
    ("EPO Patent-Import", "00BCC4", [
        ("Quelle", "195 GB Nested ZIPs\n(docdb_xml_bck_*.zip)"),
        ("Extraktion", "Outer ZIP \u2192 Inner ZIPs\n\u2192 XML iterparse()"),
        ("Parsing", "Pub-Nr, Titel, CPC-Codes\nApplicants, Filing-Date"),
        ("Staging", "TEMP TABLE + COPY\n10K Records/Batch"),
        ("Insert", "ON CONFLICT DO NOTHING\nPartition by pub_year"),
        ("Post", "Trigger: search_vector\nTrigger: publication_year"),
    ]),
    ("CORDIS Projekt-Import", "FFB800", [
        ("Quelle", "JSON-ZIP-Archive\nFP7, H2020, HORIZON"),
        ("Parsing", "JSON Array oder CSV\nDate/Amount-Konvertierung"),
        ("Validierung", "framework CHECK\nrole + activity_type"),
        ("Staging", "COPY-Protokoll\nDedup: (project_id, name)"),
        ("Insert", "Projects \u2192 Organizations\n\u2192 Publications"),
        ("Post", "EuroSciVoc-Extraktion\nHierarchie + Junction"),
    ]),
    ("API-Cache-Pipeline", "E91E63", [
        ("OpenAIRE", "Publikationstrends\n7 Tage Cache-TTL"),
        ("Semantic Scholar", "Papers + Authors\n30 Tage Cache-TTL"),
        ("GLEIF", "LEI Entity-Lookup\n90 Tage Cache-TTL"),
        ("Modus", "On-Demand per Anfrage\nStale-After Invalidierung"),
        ("Speicher", "research_schema.papers\nresearch_schema.authors"),
        ("Purge", "purge_stale_papers()\npurge_stale_gleif()"),
    ]),
]

for idx, (title, color, steps) in enumerate(pipelines):
    x = col_start_x + idx * (col_w + col_gap)
    box(s, x, col_y, col_w, col_h)
    accent_bar(s, x, col_y, col_h, color)

    txt(s, x + 150000, col_y + 100000, col_w - 200000, 350000,
        title, 15, color=color, bold=True, align="center")

    for i, (step_title, step_desc) in enumerate(steps):
        sy = col_y + 550000 + i * 800000
        txt(s, x + 180000, sy, col_w - 300000, 250000,
            step_title, 12, color="FFFFFF", bold=True)
        txt(s, x + 180000, sy + 250000, col_w - 300000, 450000,
            step_desc, 11, color="8FA8B5")

# Footer
txt(s, M, 6500000, sw - 2 * M, 300000,
    "Scheduler: APScheduler Cron (So 02:00 UTC)  |  Batch: 10K Records  |  "
    "Protokoll: asyncpg COPY  |  Tracking: cross_schema.import_log",
    11, color="78909C", align="center")

print("Added: Import-Pipeline Details")

# ================================================================
# FOLIE 3: TRANSFORMATIONEN & MATERIALIZED VIEWS
# ================================================================
s = setup_slide(prs)
txt(s, M, 150000, sw - 2 * M, 450000, "TRANSFORMATIONEN & MATERIALIZED VIEWS", 28,
    color="FFFFFF", bold=True, align="center")

# Left: Transformations
left_x = M + 200000
left_w = 5200000
left_y = 800000

transforms = [
    ("CPC-Normalisierung", "Leerzeichen entfernt, VARCHAR(8) Subclass-Level\nGIN-Index für Array-Containment (@>)",
     "00BCC4"),
    ("Applicant-Normalisierung", "Corporate-Suffixe entfernt (GmbH, Inc, Ltd)\nnormalized_name + pg_trgm Fuzzy-Index",
     "00BCC4"),
    ("Entity Resolution", "pg_trgm Similarity-Matching über EPO/CORDIS/GLEIF\n129.8K Unified Actors mit Confidence-Score",
     "FFB800"),
    ("Range-Partitioning", "patents + patent_applicants + patent_cpc nach Jahr\nCo-Partitioning für Partition-Pruning bei JOINs",
     "E91E63"),
    ("Volltextsuche", "tsvector mit gewichteten Feldern (Titel=A, CPC=B)\nTrigger-basiert, ti_plainto_tsquery() + unaccent",
     "4CAF50"),
]

for i, (title, desc, color) in enumerate(transforms):
    ty = left_y + i * 1050000
    box(s, left_x, ty, left_w, 950000)
    accent_bar(s, left_x, ty, 950000, color)
    txt(s, left_x + 180000, ty + 100000, left_w - 300000, 300000,
        title, 14, color=color, bold=True)
    txt(s, left_x + 180000, ty + 420000, left_w - 300000, 450000,
        desc, 11, color="8FA8B5")

# Right: Materialized Views
right_x = M + 5700000
right_w = 5400000

box(s, right_x, left_y, right_w, 5200000)
accent_bar(s, right_x, left_y, 5200000, "4CAF50")

txt(s, right_x + 180000, left_y + 100000, right_w - 300000, 350000,
    "9 MATERIALIZED VIEWS", 16, color="4CAF50", bold=True, align="center")
txt(s, right_x + 180000, left_y + 450000, right_w - 300000, 300000,
    "REFRESH CONCURRENTLY (non-blocking, UNIQUE INDEX)", 11,
    color="8FA8B5", align="center")

mvs = [
    ("Patent-MVs (nach EPO-Import)", "00BCC4", [
        "mv_patent_counts_by_cpc_year",
        "mv_cpc_cooccurrence (Top-200 Codes)",
        "mv_top_applicants (\u226510 Patente)",
        "mv_patent_country_distribution",
        "mv_yearly_tech_counts",
    ]),
    ("CORDIS-MVs (nach CORDIS-Import)", "FFB800", [
        "mv_project_counts_by_year",
        "mv_cordis_country_pairs",
        "mv_top_cordis_orgs (\u22653 Projekte)",
        "mv_funding_by_instrument",
    ]),
]

my = left_y + 850000
for group_title, color, views in mvs:
    txt(s, right_x + 250000, my, right_w - 400000, 300000,
        group_title, 12, color=color, bold=True)
    my += 350000
    for view in views:
        txt(s, right_x + 350000, my, right_w - 500000, 250000,
            "\u2022 " + view, 11, color="B8CAD4")
        my += 280000
    my += 150000

# Refresh functions
txt(s, right_x + 250000, my + 100000, right_w - 400000, 600000,
    "refresh_patent_views()  \u2192  5 Patent-MVs\n"
    "refresh_cordis_views()  \u2192  4 CORDIS-MVs\n"
    "refresh_all_views()       \u2192  Alle 9 MVs",
    11, color="78909C")

print("Added: Transformationen & MVs")

# ================================================================
# Move 3 new slides after Folie 12 (last Datenkonzept schema slide)
# ================================================================
sldIdLst = prs.slides._sldIdLst
all_items = list(sldIdLst)
new_slides = all_items[-3:]
old_slides = all_items[:-3]

# Insert after index 11 (Folie 12 = export_schema)
reordered = old_slides[:12] + new_slides + old_slides[12:]

for item in list(sldIdLst):
    sldIdLst.remove(item)
for item in reordered:
    sldIdLst.append(item)

prs.save(pptx_path)
print(f"\nDone. Total slides: {len(prs.slides)}")
print("3 ETL-Folien eingefügt nach Folie 12 (Schema-Block).")
