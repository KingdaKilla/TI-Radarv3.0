"""Add Herausforderungen slide to TI-Radar presentation."""
import os
import shutil
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN
from lxml import etree

BASE_DIR = r"C:\Users\bensc\OneDrive\HWR\6. Semester\Bachelor\07_Prototypen\mvp_v3.0_new"
pptx_path = os.path.join(BASE_DIR, "TI-Radar_MVP_v3.0_Praesentation.pptx")
bg_image = os.path.join(BASE_DIR, "figures", "slide_background_dark_v2.png")

# Work on a temp copy to avoid OneDrive locking issues
tmp_path = os.path.join(os.environ["TEMP"], "pres_challenges.pptx")
shutil.copy2(pptx_path, tmp_path)

prs = Presentation(tmp_path)
slide_width = prs.slide_width
slide_height = prs.slide_height


def add_text_box(slide, left, top, width, height, text, font_size,
                 color="B8CAD4", bold=False, alignment="left"):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if alignment == "center":
        p.alignment = PP_ALIGN.CENTER

    pPr = p._pPr
    if pPr is None:
        pPr = etree.SubElement(p._p, qn("a:pPr"))

    defRPr = etree.SubElement(pPr, qn("a:defRPr"))
    defRPr.set("sz", str(int(font_size * 100)))
    defRPr.set("b", "1" if bold else "0")

    solidFill = etree.SubElement(defRPr, qn("a:solidFill"))
    srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgb.set("val", color)

    latin = etree.SubElement(defRPr, qn("a:latin"))
    latin.set("typeface", "Segoe UI")

    run = p.add_run()
    run.text = text
    return txBox


def add_box(slide, left, top, width, height,
            fill_color="162138", alpha="80000", border_color="1B4D6E"):
    shape = slide.shapes.add_shape(1, Emu(left), Emu(top), Emu(width), Emu(height))
    spPr = shape._element.find(qn("p:spPr"))

    for child in list(spPr):
        if child.tag in (qn("a:solidFill"), qn("a:noFill")):
            spPr.remove(child)

    sf = etree.SubElement(spPr, qn("a:solidFill"))
    sr = etree.SubElement(sf, qn("a:srgbClr"))
    sr.set("val", fill_color)
    a = etree.SubElement(sr, qn("a:alpha"))
    a.set("val", alpha)

    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    ln.set("w", str(int(Pt(0.75))))
    for child in list(ln):
        if child.tag in (qn("a:solidFill"), qn("a:noFill")):
            ln.remove(child)
    lf = etree.SubElement(ln, qn("a:solidFill"))
    ls = etree.SubElement(lf, qn("a:srgbClr"))
    ls.set("val", border_color)
    la = etree.SubElement(ls, qn("a:alpha"))
    la.set("val", "60000")

    prstGeom = spPr.find(qn("a:prstGeom"))
    if prstGeom is not None:
        prstGeom.set("prst", "roundRect")
    return shape


def add_accent_bar(slide, x, y, height, color):
    bar = slide.shapes.add_shape(1, Emu(x), Emu(y), Emu(60000), Emu(height))
    spPr = bar._element.find(qn("p:spPr"))
    for child in list(spPr):
        if child.tag in (qn("a:solidFill"), qn("a:noFill")):
            spPr.remove(child)
    sf = etree.SubElement(spPr, qn("a:solidFill"))
    sr = etree.SubElement(sf, qn("a:srgbClr"))
    sr.set("val", color)
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    for child in list(ln):
        ln.remove(child)
    etree.SubElement(ln, qn("a:noFill"))
    return bar


# ===== Create slide =====
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# Background image
pic = slide.shapes.add_picture(bg_image, left=0, top=0,
                                width=slide_width, height=slide_height)
sp_tree = slide.shapes._spTree
pic_element = pic._element
sp_tree.remove(pic_element)
sp_tree.insert(2, pic_element)

# Remove solid bg fill
bg = slide.background
bgPr = bg._element.find(qn("p:bgPr"))
if bgPr is not None:
    solidFill = bgPr.find(qn("a:solidFill"))
    if solidFill is not None:
        bgPr.remove(solidFill)
        etree.SubElement(bgPr, qn("a:noFill"))

# ===== Title =====
LEFT_MARGIN = 457200
CONTENT_WIDTH = slide_width - LEFT_MARGIN * 2

add_text_box(slide, LEFT_MARGIN, 200000, CONTENT_WIDTH, 500000,
             "HERAUSFORDERUNGEN", 32, color="FFFFFF", bold=True,
             alignment="center")

# ===== 6 Challenge items: 2 columns x 3 rows =====
items = [
    ("Zu kleine Festplatte",
     "Server hatte nur 320 GB\n\u2192 Hardware-Aufruestung war noetig", "FF6B35"),
    ("Datenmenge unterschaetzt",
     "Importvolumen war deutlich\ngroesser als urspruenglich geplant", "FFB800"),
    ("Upload zu langsam",
     "Upload-Geschwindigkeiten\nzum Server zu gering", "E91E63"),
    ("Komplexes DB-Schema",
     "Datenbank-Schema sehr komplex\n\u2192 schwieriger Datenimport", "2196F3"),
    ("17 Services + DB",
     "Konfiguration und Orchestrierung\nvon 17 Microservices + Datenbank", "00BCC4"),
    ("Berechtigungskonzepte",
     "Zugriffs- und Rollenkonzepte\nfuer Produktiv-Deployment", "4CAF50"),
]

col_width = 5400000
col_gap = 400000
row_height = 1250000
row_gap = 180000
start_y = 900000
col1_x = LEFT_MARGIN + 200000
col2_x = col1_x + col_width + col_gap

for idx, (title, desc, accent_color) in enumerate(items):
    col = idx % 2
    row = idx // 2
    x = col1_x if col == 0 else col2_x
    y = start_y + row * (row_height + row_gap)

    add_box(slide, x, y, col_width, row_height)
    add_accent_bar(slide, x, y, row_height, accent_color)

    add_text_box(slide, x + 180000, y + 150000, col_width - 300000, 400000,
                 title, 16, color=accent_color, bold=True)
    add_text_box(slide, x + 180000, y + 550000, col_width - 300000, 650000,
                 desc, 13, color="8FA8B5")

# ===== Move slide before "Vielen Dank" (swap last two) =====
slide_list = prs.slides._sldIdLst
items_list = list(slide_list)
vielen_dank = items_list[-2]
challenges = items_list[-1]
slide_list.remove(vielen_dank)
slide_list.remove(challenges)
slide_list.append(challenges)
slide_list.append(vielen_dank)

# Save to temp, then atomically replace original
out_tmp = pptx_path + ".tmp"
prs.save(out_tmp)
os.remove(tmp_path)

try:
    os.replace(out_tmp, pptx_path)
    print(f"Done. Total slides: {len(prs.slides)}")
    print("Herausforderungen-Folie eingefuegt vor 'Vielen Dank'.")
except PermissionError:
    print(f"OneDrive sperrt die Originaldatei.")
    print(f"Neue Datei gespeichert unter: {out_tmp}")
    print("Bitte PowerPoint/OneDrive schliessen und umbenennen.")
