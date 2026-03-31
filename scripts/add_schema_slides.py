"""
Replace empty DATENKONZEPT slides (7-11) with schema zoom slides.
Each slide shows the full ERM dimmed in the background + the schema PNG zoomed in.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN
from lxml import etree
from PIL import Image
import os

pptx_path = r"C:\Users\bensc\OneDrive\HWR\6. Semester\Bachelor\07_Prototypen\mvp_v3.0_new\TI-Radar_MVP_v3.0_Praesentation.pptx"
bg_image = r"C:\Users\bensc\OneDrive\HWR\6. Semester\Bachelor\07_Prototypen\mvp_v3.0_new\figures\slide_background_dark_v2.png"
erm_full = r"C:\Users\bensc\OneDrive\HWR\6. Semester\Bachelor\07_Prototypen\mvp_v3.0_new\docs\diagrams\erm_ti_radar.png"
schema_dir = r"C:\Users\bensc\OneDrive\HWR\6. Semester\Bachelor\07_Prototypen\mvp_v3.0_new\docs\diagrams\schemas"

prs = Presentation(pptx_path)
slide_w = prs.slide_width   # 12192000 EMU
slide_h = prs.slide_height  # 6858000 EMU

# Schema order following data flow through the diagram
SCHEMAS = [
    ("patent_schema", "patent_schema (154.8M Patente, ~191 GB)"),
    ("cordis_schema", "cordis_schema (80.5K Projekte, ~1.7 GB)"),
    ("research_schema", "research_schema (Semantic Scholar Cache)"),
    ("entity_schema", "entity_schema (129.8K Unified Actors)"),
    ("cross_schema", "cross_schema (9 Materialized Views)"),
    ("export_schema", "export_schema (Analysis Cache & Reports)"),
]


def add_text_box(slide, left, top, width, height, text, font_size,
                 color="FFFFFF", bold=True, alignment="left"):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if alignment == "center":
        p.alignment = PP_ALIGN.CENTER
    elif alignment == "left":
        p.alignment = PP_ALIGN.LEFT

    pPr = p._pPr
    if pPr is None:
        pPr = etree.SubElement(p._p, qn("a:pPr"))
    defRPr = etree.SubElement(pPr, qn("a:defRPr"))
    defRPr.set("sz", str(int(font_size * 100)))
    defRPr.set("b", "1" if bold else "0")
    sf = etree.SubElement(defRPr, qn("a:solidFill"))
    sr = etree.SubElement(sf, qn("a:srgbClr"))
    sr.set("val", color)
    lat = etree.SubElement(defRPr, qn("a:latin"))
    lat.set("typeface", "Segoe UI")
    run = p.add_run()
    run.text = text
    return txBox


def set_image_opacity(pic_element, alpha_pct):
    """Set opacity on a picture element via alphaModFix."""
    blipFill = pic_element.find(qn("p:blipFill"))
    if blipFill is None:
        return
    blip = blipFill.find(qn("a:blip"))
    if blip is None:
        return
    # Add alphaModFix effect
    # Remove existing
    for ext in blip.findall(qn("a:alphaModFix")):
        blip.remove(ext)
    alpha = etree.SubElement(blip, qn("a:alphaModFix"))
    alpha.set("amt", str(int(alpha_pct * 1000)))  # 25% = 25000


# =========================================================
# Step 1: Delete slides 7-11 (indices 6-10)
# =========================================================
slide_ids_to_remove = []
sldIdLst = prs.slides._sldIdLst
all_ids = list(sldIdLst)

# Indices 6..10 are the empty DATENKONZEPT slides
for idx in range(10, 5, -1):  # reverse order: 10, 9, 8, 7, 6
    sld_id = all_ids[idx]
    rId = sld_id.get(qn("r:id"))
    # Remove from sldIdLst
    sldIdLst.remove(sld_id)
    # Remove relationship and slide part
    prs.part.drop_rel(rId)
    print(f"Removed slide at index {idx}")

# =========================================================
# Step 2: Add 6 schema zoom slides after slide 6 (index 5)
# =========================================================
blank_layout = prs.slide_layouts[6]

new_slide_ids = []
for schema_name, schema_title in SCHEMAS:
    schema_png = os.path.join(schema_dir, f"{schema_name}.png")
    if not os.path.exists(schema_png):
        print(f"WARNING: {schema_png} not found, skipping")
        continue

    slide = prs.slides.add_slide(blank_layout)

    # --- Background image (tech grid) ---
    bg_pic = slide.shapes.add_picture(bg_image, 0, 0, slide_w, slide_h)
    sp_tree = slide.shapes._spTree
    bg_elem = bg_pic._element
    sp_tree.remove(bg_elem)
    sp_tree.insert(2, bg_elem)

    # Remove solid bg fill
    bg = slide.background
    bgPr = bg._element.find(qn("p:bgPr"))
    if bgPr is not None:
        solidFill = bgPr.find(qn("a:solidFill"))
        if solidFill is not None:
            bgPr.remove(solidFill)
            etree.SubElement(bgPr, qn("a:noFill"))

    # --- Full ERM diagram dimmed as context (small, top-right corner) ---
    erm_w = Emu(2800000)   # ~7.3cm
    erm_h = Emu(2450000)   # ~6.4cm (maintain ~9824:8578 ratio)
    erm_left = slide_w - erm_w - Emu(150000)  # right margin
    erm_top = Emu(150000)   # top margin
    erm_pic = slide.shapes.add_picture(erm_full, erm_left, erm_top, erm_w, erm_h)
    # Dim the full ERM to ~20% opacity
    set_image_opacity(erm_pic._element, 20)

    # --- Schema PNG centered and large ---
    # Get aspect ratio of schema image
    with Image.open(schema_png) as pil_img:
        img_w, img_h = pil_img.size
    aspect = img_w / img_h

    # Available area for schema image
    max_w = slide_w - Emu(600000)       # margins
    max_h = slide_h - Emu(1200000)      # title + bottom margin
    top_offset = Emu(750000)            # below title

    # Fit to available area maintaining aspect ratio
    if aspect > (max_w / max_h):
        # Width-constrained
        pic_w = max_w
        pic_h = int(max_w / aspect)
    else:
        # Height-constrained
        pic_h = max_h
        pic_w = int(max_h * aspect)

    pic_left = (slide_w - pic_w) // 2
    pic_top = top_offset + (max_h - pic_h) // 2

    schema_pic = slide.shapes.add_picture(schema_png, pic_left, pic_top, pic_w, pic_h)

    # --- Title ---
    add_text_box(slide, Emu(300000), Emu(100000), Emu(4000000), Emu(450000),
                 "DATENKONZEPT", 14, color="8FA8B5", bold=False)
    add_text_box(slide, Emu(300000), Emu(350000), Emu(8000000), Emu(450000),
                 schema_title, 22, color="FFFFFF", bold=True)

    new_slide_ids.append(slide)
    print(f"Added: {schema_name}")

# =========================================================
# Step 3: Move new slides to correct position (after slide 6)
# =========================================================
# New slides are appended at the end. Move them after index 5 (slide 6).
sldIdLst = prs.slides._sldIdLst
all_items = list(sldIdLst)

# The new slides are the last 6 items
num_new = len(SCHEMAS)
new_items = all_items[-num_new:]
old_items = all_items[:-num_new]

# Insert after position 5 (slide 6 = Gesamt-ERM)
insert_after = 5
reordered = old_items[:insert_after + 1] + new_items + old_items[insert_after + 1:]

# Clear and re-add in order
for item in list(sldIdLst):
    sldIdLst.remove(item)
for item in reordered:
    sldIdLst.append(item)

prs.save(pptx_path)
print(f"\nDone. Total slides: {len(prs.slides)}")
print("Schema zoom slides inserted after Folie 6 (Gesamt-ERM).")
