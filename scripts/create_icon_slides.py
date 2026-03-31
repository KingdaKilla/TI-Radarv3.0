import sys, io
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

path = r'C:\Users\bensc\OneDrive\HWR\6. Semester\Bachelor\07_Prototypen\mvp_v3.0_new\TI-Radar_MVP_v3.0_Praesentation.pptx'
prs = Presentation(path)

bg_blob = None
for shape in prs.slides[0].shapes:
    if shape.shape_type == 13 and shape.width == 12192000 and shape.height == 6858000:
        bg_blob = shape.image.blob
        break

SW = 12192000
SH = 6858000
blank_layout = prs.slide_layouts[6]


def add_text(slide, left, top, width, height, text, size=14, bold=False, color='FFFFFF', align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor.from_string(color)
    p.alignment = align


def add_icon(slide, cx, cy, size, symbol, line_color, sym_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(cx - size // 2), Emu(cy - size // 2), Emu(size), Emu(size))
    shape.fill.background()
    shape.line.color.rgb = RGBColor.from_string(line_color)
    shape.line.width = Pt(2.5)
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = symbol
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(sym_color)
    p.alignment = PP_ALIGN.CENTER
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    bodyPr.set('anchor', 'ctr')


def add_bullets(slide, left, top, width, items, size=14, color='B0BEC5'):
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(len(items) * 370000))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor.from_string(color)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(4)


def create_slide(prs, after_idx, title, subtitle, phases):
    slide = prs.slides.add_slide(blank_layout)
    sldIdLst = prs.slides._sldIdLst
    elem = sldIdLst[-1]
    sldIdLst.remove(elem)
    sldIdLst.insert(after_idx + 1, elem)
    slide.shapes.add_picture(io.BytesIO(bg_blob), Emu(0), Emu(0), Emu(SW), Emu(SH))
    add_text(slide, 457200, 200000, 11277600, 600000, title, size=36, bold=True, color='FFFFFF', align=PP_ALIGN.LEFT)
    add_text(slide, 457200, 700000, 11277600, 350000, subtitle, size=18, color='90CAF9', align=PP_ALIGN.LEFT)
    n = len(phases)
    col_w = min(2500000, (SW - 900000) // n - 300000)
    gap = 350000
    total = n * col_w + (n - 1) * gap
    sx = (SW - total) // 2
    icon_s = 750000
    icon_cy = 1850000
    for i, ph in enumerate(phases):
        cx = sx + i * (col_w + gap) + col_w // 2
        add_icon(slide, cx, icon_cy, icon_s, ph[0], ph[1], ph[2])
        add_text(slide, cx - col_w // 2, icon_cy + icon_s // 2 + 180000, col_w, 500000,
                 ph[3], size=18, bold=True, color=ph[4])
        add_bullets(slide, cx - col_w // 2, icon_cy + icon_s // 2 + 650000, col_w,
                    ph[5], size=14, color='B0BEC5')
        if i < n - 1:
            add_text(slide, cx + col_w // 2 + 10000, icon_cy - 80000, gap - 20000, 250000,
                     '\u25B8', size=28, color='546E7A')


# INSERT FROM BACK TO FRONT

# Slide 27: HERAUSFORDERUNGEN
create_slide(prs, 26, 'HERAUSFORDERUNGEN', 'Zentrale Probleme und deren Umgang', [
    ('\u26A0', 'E91E63', 'F48FB1', 'HARDWARE', 'F48FB1', ['Festplatte zu klein', 'Upload zu langsam', 'Hardware-Upgrade']),
    ('\u2353', 'FF9800', 'FFCC80', 'DATEN', 'FFCC80', ['Menge unterschaetzt', 'Komplexes Schema', 'Schwieriger Import']),
    ('\u2699', '9C27B0', 'CE93D8', 'SERVICES', 'CE93D8', ['17 Services + DB', 'Orchestrierung', 'Konfiguration']),
    ('\u2611', '4CAF50', 'A5D6A7', 'SICHERHEIT', 'A5D6A7', ['Berechtigungen', 'Rollen-Management', 'Deployment']),
])

# Slide 19: 13 USE CASES
create_slide(prs, 18, '13 USE CASES', '4 analytische Cluster', [
    ('\u2160', '2196F3', '90CAF9', 'TECHNOLOGIE\n& REIFE', '90CAF9', ['UC1 Landscape', 'UC2 Maturity', 'UC5 CPC-Flow']),
    ('\u2161', '4CAF50', 'A5D6A7', 'MARKT-\nAKTEURE', 'A5D6A7', ['UC3 Competitive', 'UC8 Temporal', 'UC11 Actor-Type']),
    ('\u2162', '9C27B0', 'CE93D8', 'FORSCHUNG\n& FOERDERUNG', 'CE93D8', ['UC4 Funding', 'UC7 Research', 'UC10 EuroSciVoc']),
    ('\u2163', 'FF9800', 'FFCC80', 'GEOGRAPHIE\n& EXPORT', 'FFCC80', ['UC6 Geographic', 'UC9 Cross-Source', 'UC12/13 Radar']),
])

# Slide 18: TECH STACK
create_slide(prs, 17, 'TECH STACK & QUALITAET', 'Technologien im Ueberblick', [
    ('\u25A3', '00BCD4', '80DEEA', 'FRONTEND', '80DEEA', ['Next.js 14', 'TypeScript', 'TailwindCSS', 'Recharts + Nivo']),
    ('\u2630', '4CAF50', 'A5D6A7', 'BACKEND', 'A5D6A7', ['FastAPI', 'gRPC + Protobuf', 'Python 3.12', 'asyncpg']),
    ('\u2395', 'FF9800', 'FFCC80', 'DATENBANK', 'FFCC80', ['PostgreSQL 17', 'pgvector', 'pg_trgm', 'Materialized Views']),
    ('\u2692', '5C6BC0', 'C5CAE9', 'DEVOPS', 'C5CAE9', ['Docker Compose', 'GitHub Actions', '6 Pipelines', '836+ Tests']),
])

# Slide 17: FAN-OUT PATTERN
create_slide(prs, 16, 'FAN-OUT PATTERN', 'Request-Verarbeitung und Service-Architektur', [
    ('\u2192', '00BCD4', '80DEEA', '1 REQUEST', '80DEEA', ['REST/JSON Anfrage', 'Technologie + Param', 'Vom Frontend']),
    ('\u2725', '4CAF50', 'A5D6A7', 'ORCHESTRATOR', 'A5D6A7', ['asyncio.gather()', 'Parallele Ausfuehrung', 'Timeout pro UC']),
    ('\u2263', '9C27B0', 'CE93D8', '13 SERVICES', 'CE93D8', ['Jeweils eigenstaendig', 'Clean Architecture', 'gRPC']),
    ('\u2714', 'FF9800', 'FFCC80', 'ERGEBNIS', 'FFCC80', ['Aggregierte Antwort', 'Graceful Degradation', 'Teilergebnisse']),
])

# Slide 16: SECURITY
create_slide(prs, 15, 'SECURITY & INFRASTRUKTUR', 'Sicherheitskonzept in 4 Schichten', [
    ('\u2387', 'E91E63', 'F48FB1', 'VERSCHLUESSELUNG', 'F48FB1', ['LUKS Festplatte', 'HTTPS / TLS', "Let's Encrypt"]),
    ('\u2339', 'FF9800', 'FFCC80', 'NETZWERK', 'FFCC80', ['Hetzner Firewall', 'Server Firewall', 'SSH Key-Only']),
    ('\u2302', '4CAF50', 'A5D6A7', 'WEB-SECURITY', 'A5D6A7', ['Nginx Reverse Proxy', 'Security-Header', 'HSTS, CSP']),
    ('\u2610', '5C6BC0', 'C5CAE9', 'ISOLATION', 'C5CAE9', ['Docker-Container', 'Prozess-Isolation', 'Basic Auth']),
])

# Slide 15: TRANSFORMATIONEN
create_slide(prs, 14, 'TRANSFORMATIONEN', 'Datenverarbeitung nach Import', [
    ('\u2261', '2196F3', '90CAF9', 'NORMALISIERUNG', '90CAF9', ['CPC-Codes bereinigt', 'Namen vereinheitlicht', 'Duplikate entfernt']),
    ('\u2194', '9C27B0', 'CE93D8', 'VERKNUEPFUNG', 'CE93D8', ['Entity Resolution', 'Fuzzy-Matching', '129.8K Unified Actors']),
    ('\u2637', 'FF9800', 'FFCC80', 'PARTITIONIERUNG', 'FFCC80', ['Range nach Jahr', 'Co-Partitioning', 'Volltextsuche']),
    ('\u229A', '4CAF50', 'A5D6A7', 'AGGREGATION', 'A5D6A7', ['9 Materialized Views', '5 Patent-MVs', '4 CORDIS-MVs']),
])

# Slide 14: IMPORT-PIPELINE
create_slide(prs, 13, 'IMPORT-PIPELINE', 'Datenimport in 5 Schritten', [
    ('\u2B07', '78909C', 'B0BEC5', 'QUELLE', 'B0BEC5', ['EPO: 195 GB ZIPs', 'CORDIS: JSON/CSV', 'APIs: On-Demand']),
    ('\u2702', '2196F3', '90CAF9', 'EXTRAKTION', '90CAF9', ['ZIP entpacken', 'XML/JSON parsen', 'Felder extrahieren']),
    ('\u2713', '4CAF50', 'A5D6A7', 'VALIDIERUNG', 'A5D6A7', ['Typ-Konvertierung', 'Constraints', 'Deduplizierung']),
    ('\u2B06', 'FF9800', 'FFCC80', 'LADEN', 'FFCC80', ['Bulk-Insert', '10K pro Batch', 'Inkrementell']),
    ('\u2699', '9C27B0', 'CE93D8', 'NACHBEREITUNG', 'CE93D8', ['Suchindex', 'Views refresh', 'Import-Log']),
])

# Slide 5: DATENKONZEPT
create_slide(prs, 4, 'DATENKONZEPT', '4 Datenquellen, 6 Schemas, 555M+ Datensaetze', [
    ('\u2318', '2196F3', '90CAF9', 'EPO DOCDB', '90CAF9', ['154.8M Patente', 'Weltweite Daten', 'CPC-Klassifikation', '191 GB']),
    ('\u2691', '4CAF50', 'A5D6A7', 'EU CORDIS', 'A5D6A7', ['80.5K Projekte', '438K Organisationen', 'EuroSciVoc', '1.7 GB']),
    ('\u270E', '9C27B0', 'CE93D8', 'OPENAIRE', 'CE93D8', ['Open-Access-Pubs', 'Jaehrliche Trends', 'API (7d Cache)']),
    ('\u2316', 'FF9800', 'FFCC80', 'SEMANTIC\nSCHOLAR', 'FFCC80', ['Zitationsanalyse', 'h-Index, Autoren', 'API (30d Cache)']),
])

# Slide 3: MOTIVATION
create_slide(prs, 2, 'MOTIVATION & ZIELSETZUNG', 'Vom Problem zur Loesung', [
    ('\u2717', 'E91E63', 'F48FB1', 'PROBLEM', 'F48FB1', ['Fragmentierte Quellen', 'Manuelle Recherche', 'Zeitaufwaendig', 'Fehleranfaellig']),
    ('\u2192', 'FF9800', 'FFCC80', 'ANSATZ', 'FFCC80', ['Automatisierung', 'Aggregation', 'Standardisierung']),
    ('\u2714', '4CAF50', 'A5D6A7', 'TI-RADAR', 'A5D6A7', ['4+ Datenquellen', '13 Use Cases', 'Echtzeit-Dashboard', 'Export']),
])

prs.save(path)
print('9 icon-slides inserted successfully')
