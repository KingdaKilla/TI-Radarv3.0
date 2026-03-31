"""Restore 5 missing slides: Dashboard, Detail, Demo, Fazit, Vielen Dank."""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN
from lxml import etree

pptx_path = r"C:\Users\bensc\OneDrive\HWR\6. Semester\Bachelor\07_Prototypen\mvp_v3.0_new\TI-Radar_MVP_v3.0_Praesentation.pptx"
bg_image = r"C:\Users\bensc\OneDrive\HWR\6. Semester\Bachelor\07_Prototypen\mvp_v3.0_new\figures\slide_background_dark_v2.png"

prs = Presentation(pptx_path)
sw = prs.slide_width
sh = prs.slide_height
blank = prs.slide_layouts[6]


def setup_slide(prs):
    slide = prs.slides.add_slide(blank)
    pic = slide.shapes.add_picture(bg_image, 0, 0, sw, sh)
    sp = slide.shapes._spTree
    el = pic._element
    sp.remove(el)
    sp.insert(2, el)
    bg = slide.background
    bgPr = bg._element.find(qn("p:bgPr"))
    if bgPr is not None:
        sf = bgPr.find(qn("a:solidFill"))
        if sf is not None:
            bgPr.remove(sf)
            etree.SubElement(bgPr, qn("a:noFill"))
    return slide


def txt(slide, left, top, width, height, text, size, color="B8CAD4",
        bold=False, align="left"):
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    pPr = p._pPr
    if pPr is None:
        pPr = etree.SubElement(p._p, qn("a:pPr"))
    dr = etree.SubElement(pPr, qn("a:defRPr"))
    dr.set("sz", str(int(size * 100)))
    dr.set("b", "1" if bold else "0")
    sf = etree.SubElement(dr, qn("a:solidFill"))
    sr = etree.SubElement(sf, qn("a:srgbClr"))
    sr.set("val", color)
    la = etree.SubElement(dr, qn("a:latin"))
    la.set("typeface", "Segoe UI")
    p.add_run().text = text
    return tb


def box(slide, left, top, width, height, fill="162138", alpha="80000",
        border="1B4D6E"):
    shape = slide.shapes.add_shape(1, Emu(left), Emu(top), Emu(width), Emu(height))
    spPr = shape._element.find(qn("p:spPr"))
    for c in list(spPr):
        if c.tag in (qn("a:solidFill"), qn("a:noFill")):
            spPr.remove(c)
    sf = etree.SubElement(spPr, qn("a:solidFill"))
    sr = etree.SubElement(sf, qn("a:srgbClr"))
    sr.set("val", fill)
    a = etree.SubElement(sr, qn("a:alpha"))
    a.set("val", alpha)
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    ln.set("w", str(int(Pt(0.75))))
    for c in list(ln):
        if c.tag in (qn("a:solidFill"), qn("a:noFill")):
            ln.remove(c)
    lf = etree.SubElement(ln, qn("a:solidFill"))
    ls = etree.SubElement(lf, qn("a:srgbClr"))
    ls.set("val", border)
    la = etree.SubElement(ls, qn("a:alpha"))
    la.set("val", "60000")
    pg = spPr.find(qn("a:prstGeom"))
    if pg is not None:
        pg.set("prst", "roundRect")
    return shape


M = 457200  # margin

# ============================================================
# Folie 21: DASHBOARD & EXECUTIVE SUMMARY
# ============================================================
s = setup_slide(prs)
txt(s, M, 200000, sw - 2 * M, 500000, "DASHBOARD & EXECUTIVE SUMMARY", 32,
    color="FFFFFF", bold=True, align="center")

items = [
    ("Cluster-Karussell", "4 thematische Cluster\nmit Vorschau-KPIs", "00BCC4"),
    ("Executive Summary", "Patente, Projekte, Publikationen\nauf einen Blick", "FFB800"),
    ("Nachvollziehbarkeit", "Datenquellen und\nVerarbeitungszeit pro UC", "4CAF50"),
    ("Technologie-Vergleich", "Side-by-Side Analyse\nzweier Technologien", "2196F3"),
]
for idx, (title, desc, col) in enumerate(items):
    x = M + 200000 + idx * 2800000
    y = 1500000
    box(s, x, y, 2500000, 3500000)
    txt(s, x + 150000, y + 200000, 2200000, 400000, title, 18, color=col, bold=True)
    txt(s, x + 150000, y + 800000, 2200000, 2000000, desc, 14, color="8FA8B5")
print("Added: Dashboard")

# ============================================================
# Folie 22: DETAILANSICHT & VERGLEICH
# ============================================================
s = setup_slide(prs)
txt(s, M, 200000, sw - 2 * M, 500000, "DETAILANSICHT & VERGLEICH", 32,
    color="FFFFFF", bold=True, align="center")
txt(s, M, 900000, sw - 2 * M, 500000, "PROGRESSIVE DISCLOSURE", 22,
    color="00BCC4", bold=True)

lines = [
    "Cluster-Ueberblick → Tab-Panels → Detailansicht",
    "Detail: Erweiterte Metriken, Analyse-Text, Datentabelle",
    "Vergleich: Zwei Technologien Side-by-Side",
    "Export: CSV-Download und Druckfunktion",
]
box(s, M + 200000, 1500000, sw - 2 * M - 400000, 3600000)
for i, line in enumerate(lines):
    txt(s, M + 500000, 1700000 + i * 700000, sw - 2 * M - 1000000, 500000,
        line, 16, color="B8CAD4")
print("Added: Detail")

# ============================================================
# Folie 23: DEMO
# ============================================================
s = setup_slide(prs)
txt(s, M, sh // 2 - 500000, sw - 2 * M, 1000000, "Demo", 48,
    color="FFFFFF", bold=True, align="center")
print("Added: Demo")

# ============================================================
# Folie 24: FAZIT & AUSBLICK
# ============================================================
s = setup_slide(prs)
txt(s, M, 200000, sw - 2 * M, 500000, "FAZIT & AUSBLICK", 32,
    color="FFFFFF", bold=True, align="center")

# ERGEBNISSE (left)
box(s, M + 100000, 1000000, 5400000, 4800000)
txt(s, M + 300000, 1100000, 5000000, 400000, "ERGEBNISSE", 20,
    color="4CAF50", bold=True)
results = [
    "Funktionsfaehiger MVP mit 13 analytischen Use Cases",
    "Skalierbare Microservice-Architektur (gRPC + Fan-Out)",
    "Integration von 4 heterogenen Datenquellen",
    "Interaktives Dashboard mit Progressive Disclosure",
    "Nachvollziehbare Analysen mit Quellenangaben",
]
for i, r in enumerate(results):
    txt(s, M + 300000, 1600000 + i * 550000, 5000000, 400000, r, 13,
        color="B8CAD4")

# AUSBLICK (right)
box(s, M + 5700000, 1000000, 5400000, 4800000)
txt(s, M + 5900000, 1100000, 5000000, 400000, "AUSBLICK", 20,
    color="00BCC4", bold=True)
outlook = [
    "KI-gestuetzte Technologieprognosen (LLM-Integration)",
    "Erweiterung der Datenquellen (USPTO, WIPO)",
    "Kubernetes-Deployment fuer Produktionsbetrieb",
    "Benutzer-Management und gespeicherte Analysen",
    "Export-Erweiterung (PDF-Reports, API-Zugang)",
    "Echtzeit-Alerting bei Technologietrends",
]
for i, o in enumerate(outlook):
    txt(s, M + 5900000, 1600000 + i * 550000, 5000000, 400000, o, 13,
        color="B8CAD4")

# Bottom quote
txt(s, M, 6100000, sw - 2 * M, 600000,
    "TI-Radar demonstriert, dass eine umfassende, datengestuetzte Technologieanalyse "
    "als Open-Source-Loesung moeglich ist - mit modernen Architekturen und offenen "
    "Datenquellen.", 12, color="8FA8B5", align="center")
print("Added: Fazit")

# ============================================================
# Folie 25: VIELEN DANK
# ============================================================
s = setup_slide(prs)
txt(s, M, sh // 2 - 1200000, sw - 2 * M, 800000, "VIELEN DANK", 48,
    color="FFFFFF", bold=True, align="center")
txt(s, M, sh // 2 - 200000, sw - 2 * M, 500000, "Fragen & Diskussion", 24,
    color="00BCC4", bold=False, align="center")
txt(s, M, sh // 2 + 600000, sw - 2 * M, 500000,
    "TI-Radar v3.0  |  Bachelorarbeit  |  HWR Berlin  |  2026", 16,
    color="8FA8B5", bold=False, align="center")
print("Added: Vielen Dank")

prs.save(pptx_path)
print(f"\nDone. Total slides: {len(prs.slides)}")
